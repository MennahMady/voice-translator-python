from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # 1. Create the Presentation Object
    prs = Presentation()

    # Helper function to add a slide
    def add_slide(title_text, content_points):
        # Layout 1 is Title + Content
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        title = slide.shapes.title
        title.text = title_text
        
        # Set Content (Bullet points)
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()  # Clear default text
        
        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(20)  # Make text readable

    # --- SLIDE 1: Title Slide ---
    # Layout 0 is Title Slide (Centered)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Real-Time Speech-to-Speech Translation\nfor Low-Resource Arabic Dialects"
    subtitle.text = "Architecting a Dialect-Aware Transformer\n\nGraduate Machine Learning Project"

    # --- SLIDE 2: The Problem ---
    add_slide("The Problem: The 'Dialect Gap'", [
        "Standard models (like Google Translate) fail on spoken dialects.",
        "They are trained on Modern Standard Arabic (MSA), not daily speech.",
        "Example: 'How are you?'",
        "   - MSA: 'Kayfa Haluk' (Formal)",
        "   - Egyptian: 'Izzayak' (What people actually say)",
        "Latency: Cloud APIs are too slow for real-time conversation."
    ])

    # --- SLIDE 3: Current Baseline (The Prototype) ---
    add_slide("Current Baseline (MVP)", [
        "I have built a functional proof-of-concept pipeline:",
        "1. Input: Google Speech Recognition (STT)",
        "2. Translation: Google Translate API (MT)",
        "3. Output: Google TTS (Text-to-Speech)",
        "Status: Working Python app supporting English <-> MSA.",
        "Limitation: It is a 'Black Box'. I cannot fine-tune it for dialects or code-switching."
    ])

    # --- SLIDE 4: Proposed Architecture (The ML Work) ---
    add_slide("Proposed Architecture: Custom Transformer", [
        "Goal: Move from APIs to a custom Neural Machine Translation (NMT) model.",
        "Model Type: Sequence-to-Sequence (Seq2Seq) Transformer.",
        "Why Transformers?",
        "   - Handle long-range dependencies better than RNNs.",
        "   - Crucial for Arabic sentence structure (VSO vs SVO).",
        "Key Components:",
        "   - Multi-Head Attention: To map dialect words to English context.",
        "   - Custom Tokenizer: Handling informal spelling (Arabizi)."
    ])

    # --- SLIDE 5: Data Strategy ---
    add_slide("Data Strategy: The MADAR Corpus", [
        "Training Data: MADAR (Multi-Arabic Dialect Applications and Resources).",
        "Why MADAR?",
        "   - Parallel sentences for 25 different Arab city dialects.",
        "   - Aligned with English and MSA.",
        "   - Focus: Travel & Tourism domain (High utility).",
        "Preprocessing:",
        "   - Cleaning diacritics and normalizing text.",
        "   - Data Augmentation using Back-Translation."
    ])

    # --- SLIDE 6: Evaluation Metrics ---
    add_slide("Evaluation & Benchmarks", [
        "How will I measure success?",
        "Quantitative Metrics:",
        "   - BLEU Score: Comparing translation quality vs. human reference.",
        "   - WER (Word Error Rate): Measuring speech recognition accuracy.",
        "Qualitative Metrics:",
        "   - Subjective testing with native speakers.",
        "   - Comparison against Google Translate API on slang phrases."
    ])

    # --- SLIDE 7: Future Roadmap ---
    add_slide("Roadmap", [
        "Phase 1 (Complete): Baseline API-based app (English <-> MSA).",
        "Phase 2 (Current): Training Transformer on Egyptian/Levantine subsets.",
        "Phase 3 (Future):",
        "   - Model Quantization (Compression).",
        "   - Running offline on edge devices (low latency).",
        "Conclusion: Moving from 'consuming APIs' to 'architecting NMT solutions'."
    ])

    # 3. Save the file
    file_name = "Project_Presentation.pptx"
    prs.save(file_name)
    print(f"✅ Presentation saved as '{file_name}'")

if __name__ == "__main__":
    create_presentation()